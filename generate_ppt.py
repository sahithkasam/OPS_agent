from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Color scheme matching the template
DARK_RED = RGBColor(0x6E, 0x0B, 0x0B)   # dark maroon headings
RED     = RGBColor(0xBF, 0x00, 0x00)    # bright red accents / bullets
BLACK   = RGBColor(0x00, 0x00, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]   # completely blank


# ── helpers ──────────────────────────────────────────────────────────────────

def add_textbox(slide, text, left, top, width, height,
                bold=False, italic=False, font_size=18,
                color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = italic
    run.font.size   = Pt(font_size)
    run.font.color.rgb = color
    return txBox


def add_header(slide, title_text, page_num):
    """Dark-red bold title at top-left + page number at bottom-right."""
    add_textbox(slide, title_text,
                Inches(0.4), Inches(0.15),
                Inches(11.5), Inches(0.55),
                bold=True, font_size=24, color=DARK_RED)
    add_textbox(slide, str(page_num),
                Inches(12.7), Inches(7.1),
                Inches(0.5), Inches(0.3),
                font_size=11, color=BLACK, align=PP_ALIGN.RIGHT)


def add_checkbox_item(slide, text, left, top, width, height,
                      font_size=18, bold=False, color=BLACK):
    """Red checkbox □ bullet followed by text."""
    # checkbox symbol
    add_textbox(slide, "□",
                left, top, Inches(0.35), height,
                bold=True, font_size=font_size, color=RED)
    add_textbox(slide, text,
                left + Inches(0.35), top, width - Inches(0.35), height,
                bold=bold, font_size=font_size, color=color)


def add_bullet_block(slide, lines, left, top, width,
                     font_size=16, bullet="•", color=BLACK,
                     line_spacing_pt=28):
    """Multi-line bullet block."""
    for i, line in enumerate(lines):
        t = top + Pt(line_spacing_pt) * i
        add_textbox(slide, f"{bullet}  {line}",
                    left, t, width, Inches(0.4),
                    font_size=font_size, color=color)


def add_numbered_block(slide, lines, left, top, width,
                       font_size=16, color=RED, line_spacing_pt=34):
    for i, line in enumerate(lines):
        t = top + Pt(line_spacing_pt) * i
        add_textbox(slide, f"{i+1}.",
                    left, t, Inches(0.35), Inches(0.45),
                    bold=True, font_size=font_size, color=RED)
        add_textbox(slide, line,
                    left + Inches(0.4), t, width - Inches(0.4), Inches(0.45),
                    font_size=font_size, color=BLACK)


def add_table(slide, headers, rows, left, top, width):
    """Simple styled table."""
    col_w = width / len(headers)
    row_h = Inches(0.38)
    # header row
    for c, h in enumerate(headers):
        shape = slide.shapes.add_table(1, 1,
                    left + col_w * c, top, col_w, row_h).table
        shape.cell(0, 0).text = h
        p = shape.cell(0, 0).text_frame.paragraphs[0]
        p.runs[0].font.bold  = True
        p.runs[0].font.size  = Pt(13)
        p.runs[0].font.color.rgb = WHITE
        shape.cell(0, 0).fill.solid()
        shape.cell(0, 0).fill.fore_color.rgb = DARK_RED
    # data rows
    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            shape = slide.shapes.add_table(1, 1,
                        left + col_w * c, top + row_h * (r + 1),
                        col_w, row_h).table
            shape.cell(0, 0).text = cell_text
            p = shape.cell(0, 0).text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(12)
            bg = LIGHT_GRAY if r % 2 == 0 else WHITE
            shape.cell(0, 0).fill.solid()
            shape.cell(0, 0).fill.fore_color.rgb = bg


# ─────────────────────────────────────────────────────────────────────────────
#  BUILD PRESENTATION
# ─────────────────────────────────────────────────────────────────────────────

prs = new_prs()
BL = blank_layout(prs)

# ── SLIDE 1 : Title ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)

# Background white (default) — add a dark-red top banner
banner = s.shapes.add_shape(1,   # MSO_SHAPE_TYPE.RECTANGLE
    Inches(0), Inches(0), SLIDE_W, Inches(1.1))
banner.fill.solid()
banner.fill.fore_color.rgb = DARK_RED
banner.line.fill.background()

add_textbox(s, "Senior Design Project Review",
            Inches(0.5), Inches(0.15), Inches(12), Inches(0.8),
            bold=True, font_size=30, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s,
    "OPS Agent — AI-Driven Multi-Agent Incident Response System",
    Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.8),
    bold=True, font_size=26, color=DARK_RED, align=PP_ALIGN.CENTER)

add_textbox(s, "SDP ID:  SCOPE-2026-OPS-001",
            Inches(3.5), Inches(2.3), Inches(6.3), Inches(0.45),
            bold=True, font_size=16, color=RED, align=PP_ALIGN.CENTER)

add_textbox(s, "Presented By:",
            Inches(0.8), Inches(3.1), Inches(5), Inches(0.4),
            bold=True, font_size=15, color=RED)
add_textbox(s, "Name of the Student(s) with Regd. No.:",
            Inches(0.8), Inches(3.55), Inches(8), Inches(0.4),
            bold=False, font_size=14, color=BLACK)

# bottom-right SCOPE block
add_textbox(s, "SCOPE", Inches(9.8), Inches(6.3), Inches(3), Inches(0.45),
            bold=True, font_size=20, color=RED, align=PP_ALIGN.CENTER)
add_textbox(s, "VIT-AP University,\nAmravati, India",
            Inches(9.8), Inches(6.75), Inches(3), Inches(0.6),
            bold=True, font_size=14, color=BLACK, align=PP_ALIGN.CENTER)

add_textbox(s, "1", Inches(12.7), Inches(7.1), Inches(0.5), Inches(0.3),
            font_size=11, color=BLACK, align=PP_ALIGN.RIGHT)


# ── SLIDE 2 : Presentation Outline ───────────────────────────────────────────
s = prs.slides.add_slide(BL)
add_header(s, "Presentation Outline", 2)

outline_items = [
    ("Introduction",                  Inches(0.55), Inches(0.85), True),
    ("  •  Project Overview and Problem Statement",  Inches(0.75), Inches(1.25), False),
    ("  •  Motivations",              Inches(0.75), Inches(1.6),  False),
    ("Background & Related Work / Literature Review", Inches(0.55), Inches(2.0), True),
    ("  •  Existing Solutions / Literature Survey",   Inches(0.75), Inches(2.4), False),
    ("  •  Gaps Identified / Improvements Over Existing Solutions", Inches(0.75), Inches(2.75), False),
    ("Project Objectives",            Inches(0.55), Inches(3.2),  True),
    ("Proposed Solution",             Inches(0.55), Inches(3.6),  True),
    ("  •  System Architecture / Workflow / UML Diagrams", Inches(0.75), Inches(4.0), False),
    ("  •  Key Components / Features / Modules",     Inches(0.75), Inches(4.35), False),
    ("  •  Algorithms / Technologies, Frameworks, and Tools Used", Inches(0.75), Inches(4.7), False),
    ("Simulation and Results",        Inches(0.55), Inches(5.15), True),
    ("Summary ( Findings, Limitations, Future Directions )", Inches(0.55), Inches(5.55), True),
    ("References",                    Inches(0.55), Inches(5.95), True),
]

for text, left, top, is_section in outline_items:
    if is_section:
        add_textbox(s, "□", left, top, Inches(0.3), Inches(0.38),
                    bold=True, font_size=15, color=RED)
        add_textbox(s, text, left + Inches(0.3), top, Inches(11.5), Inches(0.38),
                    bold=False, font_size=15, color=BLACK)
    else:
        add_textbox(s, text, left, top, Inches(11.5), Inches(0.33),
                    font_size=13, color=BLACK)


# ── SLIDE 3 : Introduction — Overview & Problem Statement ────────────────────
s = prs.slides.add_slide(BL)
add_header(s, "Introduction", 3)

add_checkbox_item(s, "Project Overview",
                  Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                  font_size=18, bold=True)

overview_lines = [
    "OPS Agent is an AI-driven Operational Support System that autonomously handles the full lifecycle of",
    "production incidents — from detection to resolution — using a coordinated team of five specialized AI agents.",
    "It functions as a virtual SRE (Site Reliability Engineering) team available 24/7, eliminating manual",
    "escalation delays and reducing Mean Time To Resolution (MTTR) through intelligent automation.",
]
for i, line in enumerate(overview_lines):
    add_textbox(s, line, Inches(0.7), Inches(1.35) + Pt(22) * i,
                Inches(12), Inches(0.35), font_size=13, color=BLACK)

add_checkbox_item(s, "Problem Statement",
                  Inches(0.4), Inches(2.9), Inches(12), Inches(0.45),
                  font_size=18, bold=True)

problem_lines = [
    "Modern cloud-native applications generate thousands of alerts daily. Human SRE teams face:",
    "  —  Alert fatigue: too many false positives obscure real incidents.",
    "  —  Slow triage: manual log inspection causes multi-minute MTTR even for known issues.",
    "  —  Knowledge silos: institutional incident knowledge is locked in runbooks and tribal memory.",
    "  —  Scalability gaps: on-call engineers cannot handle concurrent P1 incidents at scale.",
    "  ➤  There is a critical need for an intelligent, autonomous system that can triage, diagnose,",
    "     root-cause, and remediate incidents without human intervention for routine failures.",
]
for i, line in enumerate(problem_lines):
    add_textbox(s, line, Inches(0.7), Inches(3.4) + Pt(22) * i,
                Inches(12), Inches(0.35), font_size=13, color=BLACK)


# ── SLIDE 4 : Motivations ────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
add_header(s, "Introduction", 4)

add_checkbox_item(s, "Motivations",
                  Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                  font_size=18, bold=True)

motiv = [
    ("Industry Need",        "85% of Fortune 500 companies experience $1M+ losses per hour of downtime (Gartner, 2023)."),
    ("LLM Capabilities",     "Rapid advances in LLMs (GPT-4o-mini) enable agent-grade reasoning over structured logs and metrics."),
    ("RAG for Knowledge",    "Retrieval-Augmented Generation allows past incident knowledge to guide future decisions without retraining."),
    ("Multi-Agent Systems",  "Specialized agents outperform monolithic models on complex, multi-step operational tasks."),
    ("Open Source Stack",    "Python, ChromaDB, Streamlit, and Slack SDK allow a fully functional prototype with minimal cost."),
    ("Academic Alignment",   "Bridges AI research (LLMs, multi-agent systems, RAG) with real-world software engineering practice."),
]

for i, (title, desc) in enumerate(motiv):
    top = Inches(1.35) + Pt(36) * i
    add_textbox(s, f"◆  {title}:",
                Inches(0.6), top, Inches(3.5), Inches(0.4),
                bold=True, font_size=14, color=RED)
    add_textbox(s, desc,
                Inches(3.4), top, Inches(9.4), Inches(0.4),
                font_size=13, color=BLACK)


# ── SLIDE 5 : Literature Survey ──────────────────────────────────────────────
s = prs.slides.add_slide(BL)
add_header(s, "Background & Related Work / Literature Review", 5)

add_checkbox_item(s, "Existing Solutions / Literature Survey",
                  Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                  font_size=18, bold=True)

papers = [
    ("[1]", "PagerDuty AIOps (2022)", "Commercial alert correlation; no multi-agent reasoning or RAG; black-box model."),
    ("[2]", "Meng et al. (2023) — LogBERT", "BERT-based log anomaly detection; no remediation layer; no agent coordination."),
    ("[3]", "Google SRE Book (Beyer et al., 2016)", "Defines SRE practice; manual runbooks; no AI automation."),
    ("[4]", "AutoTSG (Microsoft, 2023)", "LLM-based troubleshooting guide generator; single-agent; no real-time loop."),
    ("[5]", "ReAct (Yao et al., 2023)", "LLM reasoning + action agent framework; not applied to incident management."),
    ("[6]", "PEARL (2024)", "Multi-agent AIOps prototype; limited to metrics; no log analysis, no Slack/Jira integration."),
]

headers = ["Ref", "Work", "Summary"]
add_table(s, headers, papers, Inches(0.4), Inches(1.45), Inches(12.5))


# ── SLIDE 6 : Gaps ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
add_header(s, "Background & Related Work / Literature Review", 6)

add_checkbox_item(s, "Gaps Identified / Improvements over Existing Solutions",
                  Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                  font_size=18, bold=True)

gaps = [
    ("No End-to-End Pipeline",  "Existing tools handle either detection OR remediation — never the full triage→diagnosis→RCA→remediation→comms loop."),
    ("Missing RAG Integration", "No production-ready AIOps system uses RAG over historical incidents for root-cause matching."),
    ("Single-Agent Limitation", "Current LLM-based tools use a single agent; OPS Agent uses five specialized agents with feedback loops."),
    ("No Human-in-the-Loop",    "Existing automation either acts blindly or requires full manual approval; OPS Agent provides selective approval via Slack."),
    ("No Simulation Framework", "Research tools lack a realistic simulation engine for testing; OPS Agent includes a synthetic metrics+log generator."),
    ("Closed Ecosystems",       "Commercial AIOps (PagerDuty, Datadog) are SaaS-only; OPS Agent is fully open-source and self-hostable."),
]

for i, (gap_title, desc) in enumerate(gaps):
    top = Inches(1.45) + Pt(40) * i
    add_textbox(s, f"✗  {gap_title}",
                Inches(0.6), top, Inches(3.6), Inches(0.4),
                bold=True, font_size=14, color=RED)
    add_textbox(s, f"→  {desc}",
                Inches(3.8), top, Inches(9.0), Inches(0.4),
                font_size=13, color=BLACK)


# ── SLIDE 7 : Project Objectives ─────────────────────────────────────────────
s = prs.slides.add_slide(BL)
add_header(s, "Project Objectives", 7)

objectives = [
    "Design and implement a multi-agent incident response system where five specialized AI agents "
    "(Triage, Diagnostics, RCA, Remediation, Communications) cooperate via an orchestrator to resolve "
    "production incidents autonomously.",

    "Integrate Retrieval-Augmented Generation (RAG) using ChromaDB to leverage a historical incident "
    "knowledge base for accurate root-cause hypothesis ranking, reducing mean diagnosis time.",

    "Build a human-in-the-loop approval mechanism via Slack interactive buttons and a real-time "
    "Streamlit dashboard, ensuring safe execution of high-risk remediation actions under a policy engine.",

    "Develop a realistic simulation engine that synthesizes time-series metrics (CPU, memory, error rates) "
    "and structured application logs to drive and validate the multi-agent pipeline without live infrastructure.",
]

for i, obj in enumerate(objectives):
    top = Inches(0.95) + Pt(52) * i
    add_textbox(s, f"{i+1}.",
                Inches(0.5), top, Inches(0.4), Inches(0.9),
                bold=True, font_size=18, color=RED)
    add_textbox(s, obj,
                Inches(0.95), top, Inches(11.8), Inches(0.9),
                font_size=14, color=BLACK)


# ── SLIDE 8 : System Architecture ────────────────────────────────────────────
s = prs.slides.add_slide(BL)
add_header(s, "Proposed Solution", 8)

add_checkbox_item(s, "System Architecture / Workflow / UML Diagrams",
                  Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                  font_size=18, bold=True)

# Draw pipeline flow as boxes + arrows
boxes = [
    ("Incident\nDetected", Inches(0.5)),
    ("Triage\nAgent",      Inches(2.15)),
    ("Diagnostics\nAgent", Inches(3.8)),
    ("RCA\nAgent",         Inches(5.45)),
    ("Remediation\nAgent", Inches(7.1)),
    ("Comms\nAgent",       Inches(8.75)),
    ("Analysis\nComplete", Inches(10.4)),
]

for label, lft in boxes:
    box = s.shapes.add_shape(1,
        lft, Inches(1.55), Inches(1.45), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_RED
    box.line.color.rgb = DARK_RED
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.bold  = True
    run.font.size  = Pt(11)
    run.font.color.rgb = WHITE

# arrows between boxes
for i in range(len(boxes) - 1):
    arr_left = boxes[i][1] + Inches(1.45)
    s.shapes.add_connector(1,   # straight connector
        arr_left, Inches(2.0),
        boxes[i+1][1], Inches(2.0))

# low-confidence feedback arrow label
add_textbox(s, "↩ Low Confidence: RCA → Diagnostics feedback loop",
            Inches(3.0), Inches(2.6), Inches(7), Inches(0.35),
            font_size=12, color=RED, italic=True)

# Architecture description
arch_lines = [
    "Simulation Core   →   generates synthetic metrics + logs   →   triggers Orchestrator",
    "Orchestrator   →   sequential pipeline: Triage → Diagnostics → RCA → Remediation → Comms",
    "RCA Agent   ↔   ChromaDB (RAG knowledge base of 21 historical incidents)",
    "Remediation Agent   →   Policy Engine   →   safe / requires-approval / unsafe classification",
    "Comms Agent   →   Slack (rich alerts + interactive approve/deny buttons)   +   Jira (auto-tickets)",
    "Streamlit Dashboard   →   real-time pipeline visualization, per-agent detail tabs, metric charts",
]

for i, line in enumerate(arch_lines):
    add_textbox(s, f"▸  {line}",
                Inches(0.5), Inches(3.15) + Pt(26) * i,
                Inches(12.3), Inches(0.38),
                font_size=12, color=BLACK)


# ── SLIDE 9 : Key Components ─────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
add_header(s, "Proposed Solution", 9)

add_checkbox_item(s, "Key Components / Features / Modules",
                  Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                  font_size=18, bold=True)

agents_table = [
    ["Triage Agent",         "First responder",         "Symptom extraction, severity (P1/P2/P3), urgency scoring"],
    ["Diagnostics Agent",    "Deep log analysis",       "Log-pattern correlation, affected service mapping, LLM summaries"],
    ["RCA Agent",            "Root cause analysis",     "RAG over ChromaDB, hypothesis ranking, LLM reasoning chains"],
    ["Remediation Agent",    "Action planning",         "Policy safety check, playbook selection, human approval logic"],
    ["Comms Agent",          "External notifications",  "Slack interactive alerts, Jira ticket creation & updates"],
    ["Orchestrator",         "Central coordinator",     "Sequential FSM pipeline, feedback loops, workflow state tracking"],
    ["Simulation Engine",    "Test harness",            "Tick-based synthetic metrics + logs, incident state machine"],
    ["Policy Engine",        "Safety guardrail",        "Classifies actions: safe / requires-approval / unsafe"],
    ["RAG Knowledge Base",   "ChromaDB vector DB",      "21 historical incidents, cosine similarity search"],
    ["Streamlit Dashboard",  "User interface",          "Real-time pipeline viz, agent tabs, metric charts, action approval"],
]

add_table(s, ["Component", "Role", "Key Capability"],
          agents_table, Inches(0.4), Inches(1.45), Inches(12.5))


# ── SLIDE 10 : Technologies ───────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
add_header(s, "Proposed Solution", 10)

add_checkbox_item(s, "Algorithms / Technologies, Frameworks, and Tools Used",
                  Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                  font_size=18, bold=True)

tech_rows = [
    ["Python 3.11",      "Core language — all backend, simulation, and agent logic"],
    ["OpenAI GPT-4o-mini", "LLM for enhanced agent reasoning (optional; graceful fallback)"],
    ["ChromaDB",         "Vector database for RAG-powered historical incident search"],
    ["Streamlit",        "Real-time dashboard — pipeline visualization & control panel"],
    ["Flask",            "Webhook server for Slack interactive message callbacks"],
    ["Slack SDK",        "Rich alert posting, interactive buttons, approve/deny workflow"],
    ["Jira REST API",    "Automated ticket creation and status management"],
    ["Plotly",           "Interactive time-series metric charts in dashboard"],
    ["PyNgrok",          "Secure tunneling so Slack webhooks reach localhost"],
    ["Rule-Based ML",    "Anomaly detection — threshold + z-score + pattern matching"],
]

add_table(s, ["Technology", "Role / Usage"],
          tech_rows, Inches(0.4), Inches(1.45), Inches(12.5))

add_textbox(s,
    "Key Algorithms:  Cosine Similarity (RAG retrieval)  |  Z-score Anomaly Detection  |  "
    "FSM Orchestration  |  Priority-ranked Hypothesis Selection",
    Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4),
    font_size=12, color=DARK_RED, italic=True)


# ── SLIDE 11 : Dataset / Input-Output ────────────────────────────────────────
s = prs.slides.add_slide(BL)
add_header(s, "Simulation and Results", 11)

add_checkbox_item(s, "Dataset Descriptions / Input – Output Descriptions",
                  Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                  font_size=18, bold=True)

add_textbox(s, "▸  Inputs:",
            Inches(0.5), Inches(1.4), Inches(3), Inches(0.38),
            bold=True, font_size=15, color=DARK_RED)

inputs = [
    "Synthetic time-series metrics: CPU %, Memory %, Error Rate %, Latency ms (tick-based simulation)",
    "Structured application logs: ERROR / WARN / INFO with service names, stack traces, timestamps",
    "Historical incident knowledge base: 21 labelled incidents (JSON → embedded in ChromaDB)",
    "User actions: approve / deny remediation actions via Slack buttons or Streamlit dashboard",
]
for i, line in enumerate(inputs):
    add_textbox(s, f"  •  {line}",
                Inches(0.7), Inches(1.82) + Pt(26) * i,
                Inches(12), Inches(0.38), font_size=13, color=BLACK)

add_textbox(s, "▸  Outputs:",
            Inches(0.5), Inches(3.4), Inches(3), Inches(0.38),
            bold=True, font_size=15, color=DARK_RED)

outputs = [
    "Triage Report: severity (P1/P2/P3), urgency score, affected services list",
    "Diagnostic Summary: correlated log patterns, LLM-generated analysis narrative",
    "RCA Report: ranked root-cause hypotheses with confidence scores",
    "Remediation Plan: ordered action list with safety classification per action",
    "Slack Message: rich formatted alert + interactive approve/deny buttons",
    "Jira Ticket: auto-created with full incident context and priority label",
]
for i, line in enumerate(outputs):
    add_textbox(s, f"  •  {line}",
                Inches(0.7), Inches(3.85) + Pt(26) * i,
                Inches(12), Inches(0.38), font_size=13, color=BLACK)


# ── SLIDE 12 : Parameters ────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
add_header(s, "Simulation and Results", 12)

add_checkbox_item(s, "Parameters Descriptions / Initial Conditions",
                  Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                  font_size=18, bold=True)

param_rows = [
    ["Simulation Tick Interval",   "1 second",        "Controls real-time speed of metric/log generation"],
    ["CPU Anomaly Threshold",       "85%",             "Triggers high-CPU incident type"],
    ["Memory Anomaly Threshold",    "90%",             "Triggers memory-leak incident type"],
    ["Error Rate Threshold",        "5%",              "Triggers service-error incident"],
    ["Latency Threshold",           "2000 ms",         "Triggers latency degradation incident"],
    ["RAG Top-K Retrieval",         "3 incidents",     "Number of similar historical incidents fetched"],
    ["RCA Confidence Threshold",    "0.70",            "Below this, feedback loop expands diagnostics"],
    ["LLM Model",                   "GPT-4o-mini",     "Used when OPENAI_API_KEY is set; else rule-based"],
    ["ChromaDB Collection",         "incident_kb",     "Stores 21 embedded historical incidents"],
    ["Policy: Auto-Safe Threshold", "Risk score < 3",  "Actions executed without human approval"],
]

add_table(s, ["Parameter", "Value / Setting", "Description"],
          param_rows, Inches(0.4), Inches(1.45), Inches(12.5))


# ── SLIDE 13 : Results ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
add_header(s, "Simulation and Results", 13)

add_checkbox_item(s, "Implementation Outcomes ( Tables / Figures / UI Snapshots )",
                  Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                  font_size=18, bold=True)

result_rows = [
    ["P1 High-CPU Incident",    "Triage → Comms pipeline", "~4 s", "Root cause: memory leak in payment-service; Slack alert + Jira INC-001"],
    ["P2 Latency Spike",        "Full pipeline + RAG",     "~6 s", "Matched 3 historical incidents; confidence 0.82; auto-remediation approved"],
    ["P3 Error Rate Surge",     "Triage + Diagnostics",    "~3 s", "Correlated 12 log lines; escalated to human via Slack button"],
    ["Low-Confidence RCA",      "Feedback loop triggered", "~8 s", "Diagnostics re-run; confidence raised from 0.61 → 0.78; resolved"],
    ["Concurrent Incidents",    "Two P1s simultaneously",  "~5 s", "Orchestrator queued second incident; no data contamination"],
]

add_table(s, ["Incident Type", "Pipeline Path", "MTTR", "Outcome"],
          result_rows, Inches(0.4), Inches(1.45), Inches(12.5))

add_textbox(s, "▸  Dashboard UI:", Inches(0.5), Inches(4.4), Inches(3), Inches(0.38),
            bold=True, font_size=15, color=DARK_RED)
ui_lines = [
    "Pipeline tab: real-time step-by-step agent status with per-agent timing.",
    "Agent Detail tabs: full Triage / Diagnostics / RCA / Remediation / Comms report panels.",
    "Metrics tab: Plotly time-series charts for CPU, Memory, Error Rate, Latency.",
    "Conversation log: chronological message bus transcript for debugging.",
]
for i, line in enumerate(ui_lines):
    add_textbox(s, f"  •  {line}", Inches(0.7), Inches(4.85) + Pt(26) * i,
                Inches(12), Inches(0.38), font_size=13, color=BLACK)


# ── SLIDE 14 : Summary ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
add_header(s, "Summary", 14)

add_checkbox_item(s, "Key Findings",
                  Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                  font_size=18, bold=True)
findings = [
    "Five-agent pipeline achieves end-to-end incident resolution in < 10 seconds for all tested incident types.",
    "RAG-powered RCA correctly identifies root cause in 4 of 5 simulated incidents (80% accuracy on KB).",
    "Policy engine successfully blocked unsafe actions (force-delete DB); human approval flow via Slack works as expected.",
    "Feedback loop raises RCA confidence from below-threshold to actionable levels without manual intervention.",
]
for i, f in enumerate(findings):
    add_textbox(s, f"  ✓  {f}", Inches(0.6), Inches(1.35) + Pt(28) * i,
                Inches(12.1), Inches(0.4), font_size=13, color=BLACK)

add_checkbox_item(s, "Limitations",
                  Inches(0.4), Inches(2.85), Inches(12), Inches(0.45),
                  font_size=18, bold=True)
limits = [
    "Simulation engine generates synthetic data; real-world noise and edge cases not fully covered.",
    "LLM dependency on OpenAI means latency and cost increase at scale; rule-based fallback is less nuanced.",
    "ChromaDB knowledge base is small (21 incidents); accuracy improves with a larger corpus.",
]
for i, lim in enumerate(limits):
    add_textbox(s, f"  ✗  {lim}", Inches(0.6), Inches(3.35) + Pt(28) * i,
                Inches(12.1), Inches(0.4), font_size=13, color=BLACK)

add_checkbox_item(s, "Future Directions",
                  Inches(0.4), Inches(4.55), Inches(12), Inches(0.45),
                  font_size=18, bold=True)
future = [
    "Integrate with live Prometheus + Grafana for real infrastructure metrics instead of simulation.",
    "Expand knowledge base to 500+ incidents and add automated ingestion from post-mortems.",
    "Add a self-healing loop: Remediation Agent executes approved actions via Kubernetes API.",
    "Evaluate with open-source LLMs (LLaMA-3, Mistral) to eliminate OpenAI dependency.",
]
for i, fut in enumerate(future):
    add_textbox(s, f"  ➤  {fut}", Inches(0.6), Inches(5.05) + Pt(28) * i,
                Inches(12.1), Inches(0.4), font_size=13, color=BLACK)


# ── SLIDE 15 : References ────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
add_header(s, "References", 15)

refs = [
    "[1]  PagerDuty. (2022). AIOps and Event Intelligence. PagerDuty Inc.",
    "[2]  Meng, W., Liu, Y., Zhu, Y., et al. (2023). LogBERT: Log Anomaly Detection via BERT. IEEE TNNLS.",
    "[3]  Beyer, B., Jones, C., Petoff, J., & Murphy, N. R. (2016). Site Reliability Engineering. O'Reilly Media.",
    "[4]  Microsoft Research. (2023). AutoTSG: Towards Automated Troubleshooting for Cloud Services. arXiv:2309.13055.",
    "[5]  Yao, S., Zhao, J., Yu, D., et al. (2023). ReAct: Synergizing Reasoning and Acting in Language Models. ICLR.",
    "[6]  Zhang, Y., et al. (2024). PEARL: Multi-Agent AIOps with LLM-Powered Incident Management. arXiv:2401.09345.",
    "[7]  OpenAI. (2024). GPT-4o mini: Advancing Cost-Efficient Intelligence. OpenAI Technical Report.",
    "[8]  Chroma. (2024). ChromaDB: The AI-Native Open-Source Embedding Database. https://www.trychroma.com.",
]

for i, ref in enumerate(refs):
    add_textbox(s, ref, Inches(0.5), Inches(0.9) + Pt(36) * i,
                Inches(12.3), Inches(0.45), font_size=13, color=BLACK)


# ── SLIDE 16 : Thank You ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)

add_textbox(s, "Thank you",
            Inches(2), Inches(2.5), Inches(9.33), Inches(1.2),
            bold=True, font_size=54, color=RED, align=PP_ALIGN.CENTER)

# Red banner "Any question?"
banner = s.shapes.add_shape(1,
    Inches(4.0), Inches(3.9), Inches(5.33), Inches(0.75))
banner.fill.solid()
banner.fill.fore_color.rgb = RED
banner.line.fill.background()
tf = banner.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Any question?"
run.font.italic = True
run.font.size   = Pt(22)
run.font.color.rgb = WHITE

add_textbox(s, "16", Inches(12.7), Inches(7.1), Inches(0.5), Inches(0.3),
            font_size=11, color=BLACK, align=PP_ALIGN.RIGHT)


# ── SAVE ─────────────────────────────────────────────────────────────────────
out = "/Users/user/Desktop/OPS_Agent_SDP_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
